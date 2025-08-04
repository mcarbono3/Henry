from flask import Blueprint, request, jsonify
# Mantengo los imports comentados por si quieres volver a activarlos fácilmente
# from flask_jwt_extended import jwt_required, get_jwt_identity
from models.presentation import Presentation
from app import db
import json
import os
from werkzeug.utils import secure_filename

# Importar la librería python-pptx
from pptx import Presentation as PptxPresentation

# Asegúrate de que este import es correcto y que tienes un blueprint 'ai_bp'
# from routes.ai import generate_ai_presentation

presentations_bp = Blueprint('presentations', __name__)

# --- CONFIGURACIÓN PARA SUBIDA DE ARCHIVOS ---
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}

def allowed_file(filename):
    """
    Verifica si la extensión del archivo está permitida.
    """
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Asegúrate de que la carpeta de subida existe
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# --- Función para parsear PPTX (CORREGIDA) ---
def parse_pptx_to_json(file_path):
    """
    Parsea un archivo .pptx y extrae el título y el contenido de las diapositivas.
    Retorna un diccionario con la estructura {"slides": [...]}.
    """
    slides_data = []
    try:
        prs = PptxPresentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_content_parts = []
            title = ""
            
            # Intenta obtener el título de la diapositiva
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    # CORRECCIÓN: Reemplazamos 'has_text' por una verificación de contenido vacío
                    if text_frame.text.strip():
                        # Si es un título de marcador de posición (placeholder)
                        if shape.is_placeholder and shape.placeholder_format.type == 1: # Título
                            title = text_frame.text.strip()
                        else:
                            slide_content_parts.append(text_frame.text.strip())
            
            # Si no se encontró un título de marcador de posición, usa el primer texto o un título por defecto
            if not title and slide_content_parts:
                title = slide_content_parts[0] # Usa la primera parte del contenido como título
                slide_content_parts = slide_content_parts[1:] # Elimina el primer elemento si se usó como título
            elif not title:
                title = f"Diapositiva {i + 1}" # Título por defecto

            slides_data.append({
                "title": title,
                "content": "\n".join(slide_content_parts)
            })
    except Exception as e:
        print(f"Error al parsear el archivo PPTX: {e}")
        # En caso de error, retorna un JSON vacío para evitar fallos
        return {"slides": []}
    
    return {"slides": slides_data}


@presentations_bp.route('/', methods=['GET', 'POST'])
# @jwt_required() # <--- Deshabilitado: No se requiere token JWT
def handle_presentations():
    """
    Maneja las solicitudes GET para obtener todas las presentaciones de un usuario
    y POST para crear una nueva presentación (IA, link o archivo).
    """
    # Para GET, el user_id viene como query parameter
    if request.method == 'GET':
        user_id_str = request.args.get('user_id')
        try:
            user_id = int(user_id_str) if user_id_str else None
        except (ValueError, TypeError):
            return jsonify({'error': 'El user_id proporcionado no es un número válido'}), 400

        if user_id is None:
            return jsonify({'error': 'Falta el user_id en los parámetros de la solicitud'}), 400
        
        try:
            presentations = Presentation.query.filter_by(author_id=user_id).all()
            return jsonify([p.to_dict() for p in presentations]), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    # Para POST, el user_id viene en el cuerpo (JSON o form-data)
    elif request.method == 'POST':
        try:
            # MEJORA: Lógica más robusta para determinar el tipo de fuente y datos
            if request.content_type and 'application/json' in request.content_type:
                data = request.get_json()
                source_type = data.get('source_type', 'ai')
            else: # Asumimos 'form-data' para subidas de archivos
                data = request.form
                source_type = data.get('source_type', 'upload')

            user_id_str = data.get('user_id')
            try:
                user_id = int(user_id_str) if user_id_str else None
            except (ValueError, TypeError):
                return jsonify({'error': 'El user_id proporcionado no es un número válido'}), 400

            if user_id is None:
                return jsonify({'error': 'Falta el user_id en los datos de la solicitud'}), 400
            
            if source_type == 'link':
                link_url = data.get('source_url')
                if not link_url:
                    return jsonify({'error': 'El link de la presentación es obligatorio'}), 400
                
                new_presentation = Presentation(
                    author_id=user_id,
                    title=data.get('title', 'Presentación Externa'),
                    topic=data.get('topic', 'General'), 
                    audience=data.get('audience', 'Público general'),
                    duration=data.get('duration', 'N/A'),
                    style=data.get('style', 'professional'),
                    source_type='link',
                    source_url=link_url,
                    content_json=json.dumps({}) 
                )
            
            elif source_type == 'upload':
                if 'file' not in request.files:
                    return jsonify({'error': 'No se encontró el archivo en la solicitud'}), 400
                
                file = request.files['file']
                title = data.get('title', 'Presentación Subida')
                topic = data.get('topic', 'General') 
                audience = data.get('audience', 'Público general')
                duration = data.get('duration', 'N/A')
                style = data.get('style', 'professional')
                
                if file.filename == '':
                    return jsonify({'error': 'No se seleccionó un archivo'}), 400
                
                if file and allowed_file(file.filename):
                    filename = secure_filename(file.filename)
                    file_path = os.path.join(UPLOAD_FOLDER, filename)
                    file.save(file_path)

                    presentation_content = parse_pptx_to_json(file_path)
                    slides_count = len(presentation_content.get('slides', []))
                    
                    new_presentation = Presentation(
                        author_id=user_id,
                        title=title,
                        topic=topic,
                        audience=audience,
                        duration=duration,
                        style=style,
                        source_type='upload',
                        source_url=file_path,
                        content_json=json.dumps(presentation_content),
                        slides_count=slides_count
                    )
                else:
                    return jsonify({'error': 'Formato de archivo no permitido'}), 400

            elif source_type == 'ai':
                # La lógica para 'ai' sigue siendo la misma, ya que espera JSON
                data_ai = request.get_json()
                required_fields = ['title', 'topic', 'audience', 'duration', 'style']
                if not all(field in data_ai and data_ai[field] for field in required_fields):
                    return jsonify({'error': 'Campos obligatorios faltantes para IA'}), 400
                
                user_id_str_ai = data_ai.get('user_id')
                try:
                    user_id_ai = int(user_id_str_ai) if user_id_str_ai else None
                except (ValueError, TypeError):
                    return jsonify({'error': 'El user_id proporcionado no es un número válido'}), 400
                
                if user_id_ai is None:
                    return jsonify({'error': 'Falta el user_id en los datos de la solicitud'}), 400
                
                new_presentation = Presentation(
                    author_id=user_id_ai,
                    title=data_ai.get('title'),
                    topic=data_ai.get('topic'),
                    audience=data_ai.get('audience'),
                    duration=data_ai.get('duration'),
                    style=data_ai.get('style'),
                    source_type='ai',
                    content_json=json.dumps({'slides': []})
                )
            
            else:
                return jsonify({'error': 'Tipo de fuente no válido'}), 400

            db.session.add(new_presentation)
            db.session.commit()
            
            return jsonify(new_presentation.to_dict()), 201
        
        except Exception as e:
            db.session.rollback()
            return jsonify({'error': str(e)}), 500

# NUEVA RUTA: Obtener una presentación específica por ID (GET)
@presentations_bp.route('/<int:presentation_id>', methods=['GET'])
# @jwt_required() # <--- Deshabilitado: No se requiere token JWT
def get_presentation(presentation_id):
    """
    Obtiene una presentación específica por su ID y el user_id.
    """
    user_id_str = request.args.get('user_id')
    try:
        user_id = int(user_id_str) if user_id_str else None
    except (ValueError, TypeError):
        return jsonify({'error': 'El user_id proporcionado no es un número válido'}), 400

    if user_id is None:
        return jsonify({'error': 'Falta el user_id en los parámetros de la solicitud'}), 400

    try:
        presentation = Presentation.query.filter_by(id=presentation_id, author_id=user_id).first()

        if not presentation:
            return jsonify({'error': 'Presentación no encontrada o no pertenece al usuario'}), 404

        # Retorna la presentación, incluyendo el content_json
        return jsonify(presentation.to_dict()), 200

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# RUTA para manejar la actualización de presentaciones (PUT)
@presentations_bp.route('/<int:presentation_id>', methods=['PUT'])
# @jwt_required() # <--- Deshabilitado: No se requiere token JWT
def update_presentation(presentation_id):
    """
    Actualiza una presentación por su ID.
    Ahora no requiere autenticación JWT y recibe el user_id en la URL.
    """
    user_id_str = request.args.get('user_id')
    try:
        user_id = int(user_id_str) if user_id_str else None
    except (ValueError, TypeError):
        return jsonify({'error': 'El user_id proporcionado no es un número válido'}), 400

    if user_id is None:
        return jsonify({'error': 'Falta el user_id en los parámetros de la solicitud'}), 400

    try:
        presentation = Presentation.query.filter_by(id=presentation_id, author_id=user_id).first()

        if not presentation:
            return jsonify({'error': 'Presentación no encontrada o no pertenece al usuario'}), 404

        data = request.get_json()
        if not data:
            return jsonify({'error': 'No se proporcionaron datos para actualizar'}), 400

        presentation.title = data.get('title', presentation.title)
        presentation.topic = data.get('topic', presentation.topic)
        presentation.audience = data.get('audience', presentation.audience)
        presentation.duration = data.get('duration', presentation.duration)
        presentation.style = data.get('style', presentation.style)
        presentation.source_url = data.get('source_url', presentation.source_url)

        db.session.commit()
        return jsonify(presentation.to_dict()), 200

    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@presentations_bp.route('/<int:presentation_id>', methods=['DELETE'])
# @jwt_required() # <--- Deshabilitado: No se requiere token JWT
def delete_presentation(presentation_id):
    """
    Elimina una presentación por su ID.
    Ahora no requiere autenticación JWT y recibe el user_id en la URL.
    """
    user_id_str = request.args.get('user_id')
    try:
        user_id = int(user_id_str) if user_id_str else None
    except (ValueError, TypeError):
        return jsonify({'error': 'El user_id proporcionado no es un número válido'}), 400

    if user_id is None:
        return jsonify({'error': 'Falta el user_id en los parámetros de la solicitud'}), 400
    
    try:
        presentation = Presentation.query.filter_by(id=presentation_id, author_id=user_id).first()

        if not presentation:
            return jsonify({'error': 'Presentación no encontrada o no pertenece al usuario'}), 404

        db.session.delete(presentation)
        db.session.commit()

        return jsonify({'message': 'Presentación eliminada exitosamente'}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500
